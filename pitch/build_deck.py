#!/usr/bin/env python3
"""Generate MANNMITRA.pptx pitch deck — full feature set."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

TEAL = RGBColor(0x0F, 0x9B, 0x86)
TEAL_D = RGBColor(0x0B, 0x7D, 0x6C)
SAND = RGBColor(0xFD, 0xF8, 0xF1)
INK = RGBColor(0x1F, 0x2A, 0x2E)
MUTED = RGBColor(0x5B, 0x6B, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DANGER = RGBColor(0xD6, 0x45, 0x5B)
WARN = RGBColor(0xE0, 0x89, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp
def txt(s, x, y, w, h, text, size=18, color=INK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tb
def bullets(s, x, y, w, h, items, size=16, color=INK, gap=6):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, (t_, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(gap)
        r = p.add_run(); r.text = ("\u2022 " if lvl == 0 else "\u2013 ") + t_
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb
def bg(s, color): rect(s, 0, 0, SW, SH, color)
def title_bar(s, title, kicker=None):
    bg(s, WHITE)
    rect(s, 0, 0, SW, Inches(1.15), TEAL)
    txt(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.8), title,
        size=30, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        txt(s, Inches(0.5), Inches(1.25), Inches(12), Inches(0.4), kicker,
            size=14, color=MUTED, bold=True)

# 1 TITLE
s = slide(); bg(s, TEAL)
rect(s, 0, Inches(5.4), SW, Inches(2.1), TEAL_D)
txt(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.2), "MANNMITRA", size=64, color=WHITE, bold=True)
txt(s, Inches(0.8), Inches(2.9), Inches(11.7), Inches(0.6), "\u092e\u0928 \u092e\u093f\u0924\u094d\u0930 — a friend for your mind", size=22, color=SAND)
txt(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(1.4),
    "A trust-first, tiered mental-health platform for college students.\nThe first step to feeling better shouldn't feel heavy.", size=18, color=WHITE)
txt(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.5),
    "AI \u2192 Peer \u2192 Counselor \u2192 Crisis   |   Anonymous \u00b7 No signup \u00b7 24\u00d77   |   EN / HI / MR / GU", size=15, color=WHITE, bold=True)
txt(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4), "College Hackathon Round \u2014 Pitch Deck", size=13, color=SAND)

# 2 PROBLEM
s = slide(); title_bar(s, "The Problem", "WHY EXISTING SYSTEMS FAIL")
bullets(s, Inches(0.6), Inches(1.9), Inches(6.0), Inches(4.8), [
    ("Step 1 today is 'book a counselor' \u2014 feels heavy, like confessing.", 0),
    ("Students delay or never seek help because of stigma.", 0),
    ("~15% of Indian students show depression symptoms (NIMHANS).", 0),
    ("Campus counselor shortage is real \u2014 can't be hired away.", 0),
    ("One-size-fits-all portals ignore the trust gap.", 0),
], size=18)
for i, (num, lab) in enumerate([("~15%","show depression symptoms"),("1 in 4","delay help due to stigma"),("Few","campus counselors per 1000")]):
    y = Inches(1.9 + i*1.45)
    rect(s, Inches(7.0), y, Inches(5.7), Inches(1.25), TEAL if i==0 else SAND)
    txt(s, Inches(7.2), y+Inches(0.12), Inches(2.2), Inches(1.0), num, size=30, color=WHITE if i==0 else TEAL_D, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(9.4), y+Inches(0.12), Inches(3.2), Inches(1.0), lab, size=15, color=WHITE if i==0 else INK, anchor=MSO_ANCHOR.MIDDLE)

# 3 INSIGHT
s = slide(); bg(s, SAND)
rect(s, 0, 0, SW, SH, SAND)
rect(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.12), TEAL)
txt(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.2),
    "Insight: Stigma isn't a counselor-shortage problem.\nIt's a FIRST-STEP problem.", size=34, color=TEAL_D, bold=True)
txt(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.6),
    "Lower the first step to a 2-minute anonymous mood check-in \u2014 no login, no stranger.\nEverything else (peer, counselor, crisis) builds from that moment of trust.", size=18, color=INK)
rect(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.1), TEAL)
txt(s, Inches(1.0), Inches(5.75), Inches(11.3), Inches(0.8),
    "Heavy first step  \u2192  nobody comes.   Light first step  \u2192  trust  \u2192  help scales.", size=18, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# 4 FOUR LAYERS
s = slide(); title_bar(s, "The Solution \u2014 4 Trust Layers", "NOBODY HITS A DEAD END")
layers = [
    ("1 \u00b7 AI Companion", "Anonymous 24/7. Mood check-ins, journaling, breathing. Detects severity, routes.", TEAL),
    ("2 \u00b7 Peer Support", "Trained senior students ('Peer Listeners', NIMHANS/iCall module). Anonymous rooms.", TEAL_D),
    ("3 \u00b7 Counselor", "Book with only an anonymous ID to admin. Shared pool across colleges fixes shortage.", WARN),
    ("4 \u00b7 Crisis", "One-tap helplines (iCall, Vandrevala, Tele-MANAS 14416). Auto-triggered on risk.", DANGER),
]
for i, (h, d, c) in enumerate(layers):
    x = Inches(0.6 + i*3.15)
    rect(s, x, Inches(2.0), Inches(2.95), Inches(4.4), c)
    txt(s, x+Inches(0.15), Inches(2.2), Inches(2.65), Inches(1.0), h, size=17, color=WHITE, bold=True)
    txt(s, x+Inches(0.15), Inches(3.3), Inches(2.65), Inches(3.0), d, size=14, color=WHITE)
    if i < 3:
        txt(s, x+Inches(2.95), Inches(4.0), Inches(0.25), Inches(0.5), "\u2192", size=20, color=MUTED, bold=True, align=PP_ALIGN.CENTER)

# 5 FULL FEATURE MAP
s = slide(); title_bar(s, "Full Product \u2014 What's Inside", "ONE PROTOTYPE, MANY SURFACES")
feats = [
    ("\U0001F916 AI triage", "mood + 3 Q's \u2192 severity & routing (LLM-ready)"),
    ("\U0001F91D Peer support", "anonymous rooms + 1:1 listener + chat"),
    ("\U0001F3A8 Dashboard", "mood chart, streak, badges, insights"),
    ("\U0001F4DD Journal & \U0001FAC2 CBT", "private entries + thought record"),
    ("\U0001F9ED Breathing", "animated 4-4-4-4 box breathing"),
    ("\U0001F4AC Community", "anonymous feed + kind reactions"),
    ("\U0001F6A8 SOS button", "hold-to-call \u2192 crisis overlay"),
    ("\U0001F9D1 Counselor + Admin", "shared pool, anon ledger, impact"),
    ("\U0001F30D i18n EN/HI/MR/GU", "language switch across app"),
    ("\U0001F319 Dark + Contrast", "themes persisted locally"),
    ("\U0001F6E1 Privacy + My Data", "delete-all, admin-blind"),
    ("\u2699 Gamification", "streaks & achievement badges"),
]
col_w = Inches(6.1); x0 = Inches(0.5); y0 = Inches(1.85); rh = Inches(0.78)
for i, (t_, d) in enumerate(feats):
    col = i // 6; row = i % 6
    x = x0 + col*(col_w + Inches(0.3)); y = y0 + row*rh
    rect(s, x, y, col_w, Inches(0.66), SAND)
    txt(s, x+Inches(0.12), y+Inches(0.04), col_w-Inches(0.2), Inches(0.6), f"{t_}  \u2014  {d}", size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# 6 DEMO FLOW
s = slide(); title_bar(s, "Live Demo Flow", "60-SECOND CLICKABLE STORY")
bullets(s, Inches(0.7), Inches(2.0), Inches(7.2), Inches(5.0), [
    ("Landing \u2192 'Start anonymously' (no signup, onboarding tour)", 0),
    ("Mood slider + 3 questions \u2192 rule-based triage", 0),
    ("Low/medium \u2192 routed to Peer Support room", 0),
    ("Type a risk phrase \u2192 Crisis overlay auto-triggers", 0),
    ("SOS hold-to-call button \u2192 instant helplines", 0),
    ("Dashboard: streak, badges, 'lowest on Monday' insight", 0),
    ("CBT + Community + Counselor/Admin + Privacy 'My Data'", 0),
], size=16, gap=11)
rect(s, Inches(8.4), Inches(2.0), Inches(4.3), Inches(4.6), SAND)
txt(s, Inches(8.6), Inches(2.2), Inches(3.9), Inches(0.5), "Try it live:", size=16, color=TEAL_D, bold=True)
txt(s, Inches(8.6), Inches(2.8), Inches(3.9), Inches(3.6),
    "Open index.html \u2192 checkin.html\n\ntriage.js is a pure function:\n\nscore({mood, text, answers})\n  \u2192 severity: low | medium | high\n  \u2192 route: self | peer | crisis\n\nSame interface a real LLM slots into later.", size=12, color=INK)

# 7 ARCHITECTURE
s = slide(); title_bar(s, "Architecture", "FROM DEMO TO PRODUCTION")
bullets(s, Inches(0.6), Inches(1.9), Inches(6.0), Inches(5.0), [
    ("Client: HTML/JS + global Crisis/SOS overlay", 0),
    ("Anonymous session store \u2014 zero PII (opaque anon_id)", 0),
    ("4 services: Triage \u00b7 Peer Match \u00b7 Booking \u00b7 Crisis", 0),
    ("Demo: rule-based triage + localStorage", 1),
    ("Prod: FastAPI + LLM (Gemini) behind score()", 1),
    ("Data: hashed logs, TTL auto-delete, encrypted", 0),
    ("Helpline registry: iCall / Vandrevala / Tele-MANAS", 0),
], size=15, gap=8)
rect(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.4), SAND)
txt(s, Inches(7.2), Inches(2.15), Inches(5.3), Inches(0.4), "Trust Funnel", size=16, color=TEAL_D, bold=True)
for i, l in enumerate(["Browser (anon)", "API / BFF", "Triage \u00b7 Peer \u00b7 Booking \u00b7 Crisis", "Encrypted Data Layer"]):
    y = Inches(2.7 + i*0.95)
    rect(s, Inches(7.3), y, Inches(5.1), Inches(0.75), TEAL if i%2==0 else TEAL_D)
    txt(s, Inches(7.5), y, Inches(4.7), Inches(0.75), l, size=14, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, bold=True)

# 8 PRIVACY & SAFETY
s = slide(); title_bar(s, "Privacy & Safety by Design", "JUDGES ASK \u2014 WE SHOW")
cards = [
    ("Anonymous ID", "No name/email/phone. Opaque anon_id only."),
    ("Blind booking", "Admin sees ID + slot. No faculty/parent view."),
    ("Local-first demo", "Check-in state stays on device."),
    ("Auto-purge", "Hashed logs, short TTL in production."),
    ("Consent-gated", "Identities revealed only by you."),
    ("Encrypted", "TLS + column-level encryption at rest."),
    ("\U0001F6A8 SOS + Crisis", "High-risk language auto-escalates, never routes around."),
    ("\U0001F30D Multilingual", "Stigma is worse vernacular \u2014 EN/HI/MR/GU."),
]
for i, (h, d) in enumerate(cards):
    col = i % 4; row = i // 4
    x = Inches(0.5 + col*3.12); y = Inches(2.0 + row*2.2)
    rect(s, x, y, Inches(2.92), Inches(1.95), SAND)
    txt(s, x+Inches(0.2), y+Inches(0.15), Inches(2.6), Inches(0.5), h, size=15, color=TEAL_D, bold=True)
    txt(s, x+Inches(0.2), y+Inches(0.7), Inches(2.6), Inches(1.1), d, size=12, color=INK)

# 9 MULTILINGUAL & A11Y
s = slide(); title_bar(s, "Built for Bharat", "MULTILINGUAL + ACCESSIBLE")
bullets(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.6), [
    ("Language switcher: English, Hindi, Marathi, Gujarati", 0),
    ("Resource library published in all four languages", 0),
    ("Crisis keywords detected across languages in triage", 0),
    ("Stigma is worse in vernacular, non-metro contexts", 0),
    ("High-contrast mode + keyboard focus + ARIA labels", 0),
    ("Dark mode for late-night, low-stimulation use", 0),
], size=17, gap=12)
rect(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.2), TEAL)
txt(s, Inches(7.3), Inches(2.3), Inches(5.1), Inches(0.5), "Languages", size=18, color=WHITE, bold=True)
bullets(s, Inches(7.3), Inches(3.0), Inches(5.1), Inches(3.0), [
    ("English  \u2014 EN", 0), ("Hindi  \u2014 \u0939\u093f\u0928\u094d\u0926\u0940  (HI)", 0),
    ("Marathi  \u2014 \u092e\u0930\u093e\u0920\u0940  (MR)", 0), ("Gujarati  \u2014 \u0a97\u0ac1\u0a9c\u0ab0\u0abe\u0aa4\u0ac0  (GU)", 0),
], size=17, color=WHITE, gap=12)

# 10 SCALABILITY / BUSINESS
s = slide(); title_bar(s, "Scalability & Business", "B2B2C MODEL")
bullets(s, Inches(0.6), Inches(1.9), Inches(6.1), Inches(5.0), [
    ("Sell to colleges as a subscription", 0),
    ("Cheaper than hiring 5 in-house counselors each", 1),
    ("Shared counselor pool across colleges", 0),
    ("One network serves many campuses \u2192 solves shortage", 1),
    ("Peer volunteers self-replenish each batch", 0),
    ("Near-free, compounding supply of listeners", 1),
    ("CSR / wellness-fund subsidies for low-income camps", 0),
], size=17, gap=10)
rect(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.4), TEAL)
txt(s, Inches(7.2), Inches(2.2), Inches(5.3), Inches(0.5), "Unit economics", size=18, color=WHITE, bold=True)
bullets(s, Inches(7.2), Inches(2.9), Inches(5.3), Inches(3.4), [
    ("College pays per-student/yr subscription", 0),
    ("Counselor cost amortized across N colleges", 0),
    ("Peer layer = ~0 marginal cost", 0),
    ("Retention via trust, not forced signup", 0),
], size=15, color=WHITE, gap=10)

# 11 IMPACT / METRICS
s = slide(); title_bar(s, "Impact (Demo Metrics)", "WHAT SUCCESS LOOKS LIKE")
stats = [("14","Colleges onboarded"),("38,200","Lives reached"),("1,284","Check-ins (30d)"),
         ("612","Routed to peer"),("298","Routed to counselor"),("41","Crisis escalations")]
for i, (num, lab) in enumerate(stats):
    col = i % 3; row = i // 3
    x = Inches(0.7 + col*4.1); y = Inches(2.0 + row*2.2)
    rect(s, x, y, Inches(3.7), Inches(1.9), TEAL if (i%2==0) else TEAL_D)
    txt(s, x, y+Inches(0.2), Inches(3.7), Inches(0.9), num, size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, y+Inches(1.1), Inches(3.7), Inches(0.7), lab, size=14, color=WHITE, align=PP_ALIGN.CENTER)

# 12 CREDIBILITY
s = slide(); title_bar(s, "Why It's Credible", "FRAMEWORKS & TRUST ANCHORS")
bullets(s, Inches(0.6), Inches(2.0), Inches(11.5), Inches(4.8), [
    ("Peer Listener certification aligned with NIMHANS & iCall training frameworks", 0),
    ("Crisis routing uses verified national helplines: Tele-MANAS 14416, iCall (TISS), Vandrevala", 0),
    ("Anonymous-first mirrors global best practice (7 Cups, Crisis Text Line) adapted for India", 0),
    ("Resource library in Hindi, English, Marathi & Gujarati \u2014 targets worse vernacular stigma", 0),
    ("Safety built-in: high-risk language auto-escalates, never routes around crisis", 0),
], size=18, gap=14)

# 13 ASK
s = slide(); bg(s, TEAL)
txt(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8), "The Ask", size=34, color=WHITE, bold=True)
bullets(s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(2.6), [
    ("Pilot with 3 colleges next semester \u2014 measure check-in \u2192 counselor conversion", 0),
    ("Partnership intro to a telehealth/counselor network for the shared pool", 0),
    ("Mentorship on campus wellness-fund / CSR funding path", 0),
], size=18, color=WHITE, gap=12)
rect(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.6), TEAL_D)
txt(s, Inches(1.0), Inches(4.95), Inches(11.3), Inches(1.3),
    "MANNMITRA \u2014 \u092e\u0928 \u092e\u093f\u0924\u094d\u200d\u0930. A friend for your mind.\nLight first step. Real trust. No dead ends.", size=20, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
    "Demo: mannmitra/index.html   \u00b7   In crisis: 14416 (Tele-MANAS)", size=14, color=SAND)

prs.save(r"C:\Users\Rudra\Desktop\mannmitra\pitch\MANNMITRA.pptx")
print("Saved MANNMITRA.pptx with", len(prs.slides._sldIdLst), "slides")
